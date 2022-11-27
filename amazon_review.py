from selenium import webdriver
import chromedriver_binary
import time
from time import sleep
from selenium.webdriver.chrome.options import Options
from bs4 import BeautifulSoup
import textwrap
import csv
import datetime
import pptx
import urllib.error
import urllib.request
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from pptx.dml.color import RGBColor
import shutil
import math
import re
 

# get page HTML
def get_amazon_page_info(url):
    text = ""                               #　初期化
    options = Options()                     #　オプションを用意
    options.add_argument('--incognito')     #　シークレットモードの設定を付与
    options.add_argument('--headless')      #  ヘッドレスモードのオンオフ↓
    #　chromedriverのパスとパラメータを設定
    driver = webdriver.Chrome(options=options)
    driver.get(url)                         #　chromeブラウザでurlを開く
    driver.implicitly_wait(10)              #　指定したドライバの要素が見つかるまでの待ち時間を設定
    text = driver.page_source               #　ページ情報を取得
    driver.quit()                           #　chromeブラウザを閉じる
    return text                             #　取得したページ情報を返す

# get product_status
def get_status(url):
    status_list = []
    text = get_amazon_page_info(url)    #　amazonの商品ページ情報(HTML)を取得する
    amazon_bs = BeautifulSoup(text, features='lxml')    #　HTML情報を解析する

    product_title = amazon_bs.select('#productTitle')
    product_by = amazon_bs.select('#bylineInfo')
    image_url = amazon_bs.select('.a-dynamic-image')[0].get('src')
    price = amazon_bs.select("[data-a-color='price'] > span[aria-hidden]")
    #price = amazon_bs.select(".a-color-price")
    product_description = amazon_bs.select('#feature-bullets')

    status_list.append(product_title)
    status_list.append(product_by)
    status_list.append(image_url)
    status_list.append(price)
    status_list.append(product_description)

    return status_list
    
# get product_review
def get_all_reviews(url):
    text = get_amazon_page_info(url)    #　amazonの商品ページ情報(HTML)を取得する
    amazon_bs = BeautifulSoup(text, features='lxml')    #　HTML情報を解析する

    star_list = []
    starrating = amazon_bs.select('.averageStarRatingIconAndCount')
    startable = amazon_bs.select('.histogram')

    star_list.append(starrating)
    star_list.append(startable)

    review_list = []  
    reviews = amazon_bs.select('.review') 
    # リスト内放棄
    for review in reviews:                              #　取得したレビュー数分だけ処理を繰り返す
        content_list = []
        review_title = review.select('.review-title')
        review_rating = review.select('.review-rating')
        review_text = review.select('.review-text')
        content_list.append(review_title)                      #　レビュー情報をreview_listに格納
        content_list.append(review_rating)
        content_list.append(review_text) 
        review_list.append(content_list)
        # そんなにレビュー数いらないからコメントアウトしている
        # next_page = amazon_bs.select('li.a-last a')         # 「次へ」ボタンの遷移先取得
        # # 次のページが存在する場合
        # if next_page != []: 
        #     # 次のページのURLを生成   
        #     next_url = 'https://www.amazon.co.jp/' + next_page[0].attrs['href']    
        #     url = next_url  # 次のページのURLをセットする
        #     sleep(1)        # 最低でも1秒は間隔をあける(サーバへ負担がかからないようにする)
        # else:               # 次のページが存在しない場合は処理を終了
        #     break
    return star_list,review_list

# get summarize page text
def get_page(review):
    JS_ADD_TEXT_TO_INPUT = """
    var elm = arguments[0], txt = arguments[1];
    elm.value += txt;
    elm.dispatchEvent(new Event('change'));
    """
    url = 'https://text-summary.userlocal.jp/'
    text = ""                              
    options = Options()                     
    options.add_argument('--incognito')     
    #options.add_argument('--headless')      
    driver = webdriver.Chrome(options=options)
    driver.get(url)                         
    driver.implicitly_wait(10)              
    search = driver.find_element_by_id('exampleFormControlTextarea1') 
    driver.execute_script(JS_ADD_TEXT_TO_INPUT, search, review)
    driver.implicitly_wait(10) 
    search.submit()
    # selector = '{{ CSSセレクタ }}'
    # element = WebDriverWait(driver, 30).until(
    # 	EC.visibility_of_element_located((By.CSS_SELECTOR, selector))
    # )
    text = driver.page_source   
    driver.implicitly_wait(10)           
    driver.quit()                           
    return text   

# get summarize review 
def summarize(review):
    text = get_page(review)   
    amazon_bs = BeautifulSoup(text, features='lxml')   
    sum_review = amazon_bs.select('#summaryTextarea')[0].text
    print(sum_review)
    time.sleep(1)
    sum_review = textwrap.fill(sum_review, 102).strip()
    return sum_review
 
def write_csv(url,review_url):
    status_list = get_status(url)
    get_review_list = get_all_reviews(review_url) #タプルで帰ってくる
    star_list = get_review_list[0]
    review_list = get_review_list[1]
    #CSVにレビュー情報の書き出し
    date_now = datetime.datetime.now().strftime('%Y-%m-%d %H-%M-%S')
    
    with open(f'csv/{date_now}.csv','w',encoding='utf-8_sig') as f:
        writer = csv.writer(f, lineterminator='\n')
        review_header = ['商品url','商品名','製造社',' #Shorts','【レビュー】【要約】','【review】【comparison】']
        writer.writerow(review_header)
        csv_review = []
        product_title = textwrap.fill(status_list[0][0].text, 80)
        product_by = textwrap.fill(status_list[1][0].text, 80)
        image_url = status_list[2]
        price = textwrap.fill(status_list[3][0].text, 80)
        product_description = textwrap.fill(status_list[4][0].text, 80)
        csv_review.append(url)
        csv_review.append(product_title.strip())
        csv_review.append(product_by.strip())
        csv_review.append(price.strip())
        csv_review.append(product_description.strip())
        # 概要欄用テキスト
        csv_review.append('\n'.join(csv_review))
        csv_review.append(image_url)
        writer.writerow(csv_review)
        star_header = ['average','table']
        writer.writerow(star_header)
        csv_star=[]
        starrating = textwrap.fill(star_list[0][0].text, 80)
        startable = textwrap.fill(star_list[1][0].text, 80)
        csv_star.append(starrating.strip())
        csv_star.append(startable.strip())
        writer.writerow(csv_star)
        # 全データを表示
        for i in range(len(review_list)):
            odd_csvlist=[]
            review_title = textwrap.fill(review_list[i][0][0].text, 80)  # review_list=[[title],[date],[comment]] title=['...']
            review_rating = textwrap.fill(review_list[i][1][0].text, 80)
            #review_date = textwrap.fill(review_list[i][2][0].text, 80)
            #データ作成
            odd_csvlist.append(review_rating.strip()) #レビューテキストの先頭・末尾の空白文字を除去
            #csvlist.append(review_date.strip())
            odd_csvlist.append(review_title.strip())        
            # for a in range(len(review_list[0])):
            #     review_elem = textwrap.fill(review_list[i][a][0].text, 80)
            #     csvlist.append(review_elem.strip())
            # 出力    
            writer.writerow(odd_csvlist) 
            # starの偶数番目
            even_csvlist=[]
            review = textwrap.fill(review_list[i][2][0].text, 102).strip()
            print(review)
            sum_review = summarize(review)
            even_csvlist.append(sum_review.replace('\n', ' ')) 
            writer.writerow(even_csvlist)               
        # ファイルクローズ
        f.close()
    return date_now

# download img from web
def download_file(url, dst_path):
    try:
        with urllib.request.urlopen(url) as web_file:
            data = web_file.read()
            with open(dst_path, mode='wb') as local_file:
                local_file.write(data)
    except urllib.error.URLError as e:
        print(e)

# write to pptx
def write_pptx(country,date_now):
    # create pptx file
    csv_path=f'csv/{date_now}.csv'
    output_path=f'pptx/{date_now}.pptx'
    img_path='img/'
    if country == 'USA':
        src = 'pptx/sample_usa.pptx'
    elif country == 'JAPAN':
        src = 'pptx/sample.pptx'
    copy = output_path
    shutil.copyfile(src,copy)

    prs = pptx.Presentation(output_path)
    # reccet pre_img
    res_pic = prs.slides[0].shapes[0] # 縦幅で合わせた画像は削除する
    res_pic = res_pic._element
    res_pic.getparent().remove(res_pic)
    with open(csv_path,'r',encoding='utf-8_sig') as f:
        reader = csv.reader(f)
        sample_data = [row for row in reader]
    # create data of graph
    num = re.findall('[0-9]+',sample_data[3][1])
    star_list = list(map(float, num)) # 文字列配列を数値配列に
    star_list = star_list[1::2] #奇数要素を取得
    for i in range(len(star_list)):
        star_list[i]=star_list[i]/100
    # write graph from csv to pptx
    chart = prs.slides[0].shapes[1].chart  
    chart_data = CategoryChartData()
    if country == 'USA':
        chart_data.categories = '5 star','4 star','3 star','2 star','1 star'
    elif country == 'JAPAN':
        chart_data.categories = '星5つ','星4つ','星3つ','星2つ','星1つ'    
    chart_data.add_series('系列 1', star_list)
    chart.replace_data(chart_data)
    #　テーブルの書き換え
    tbl = prs.slides[0].shapes[2].table 
    # write table from csv to pptx
    font_size = Pt(20)
    title_font_size = Pt(28)
    font_color = RGBColor(0xFF, 0xFF, 0xFF)
    line_spacing = 1.4
    with open(f'csv/{date_now}.csv','r',encoding='utf-8_sig') as f:
        reader = csv.reader(f)
        sample_data = [row for row in reader]
    for i in range(len(sample_data)):
        if i < 4 :
            pass
        else:
            # 偶数行目に白文字
            for j in range(len(sample_data[i])): 
               tbl.cell(i-4, j).text =str(sample_data[i][j])
               paras = tbl.cell(i-4, j).text_frame.paragraphs
               for para in paras:
                   para.font.size = font_size
                   if j % 2 == 1:
                       para.font.size = title_font_size
                       para.font.bold = True
                   if i % 2 == 1:
                       para.font.color.rgb = font_color
                       para.line_spacing = line_spacing
    img_url = sample_data[1][6]
    url = img_url
    dst_path = f'img/{date_now}.png'
    download_file(url, dst_path)
    width = prs.slide_width
    height = prs.slide_height
    # 画像の挿入
    pic = prs.slides[0].shapes.add_picture(dst_path, 0, 0, width = None, height = height)
    # 画像の横幅がスライドより大きい場合は横を合わせる
    if width - pic.width < 0 :
        del_pic = prs.slides[0].shapes[3] # 縦幅で合わせた画像は削除する
        del_pic = del_pic._element
        del_pic.getparent().remove(del_pic)
        pic = prs.slides[0].shapes.add_picture(dst_path, 0, 0, width = width, height = None)
    # 中心に移動
    pic.left = int( ( width  - pic.width  ) / 2 )
    pic.top  = int( ( height - pic.height ) / 2 )
    # 重なりの順序を変更
    picture = prs.slides[0].shapes[0]._element 
    picture.addprevious(prs.slides[0].shapes[3]._element) # shape[2]をshape[0]の前つまり最背面に移動

    prs.save(f'pptx/{date_now}.pptx')

def create_short(date_now):
    csv_path=f'csv/{date_now}.csv'
    output_path=f'pptx/short_{date_now}.pptx'
    img_path='img/'
    src = 'pptx/short_develop.pptx'
    copy = output_path
    shutil.copyfile(src,copy)

    prs = pptx.Presentation(output_path)
    # 前回の写真を削除してリセットする
    res_pic = prs.slides[0].shapes[1] 
    res_pic = res_pic._element
    res_pic.getparent().remove(res_pic)

    # スライド内のオブジェクトの読み取り
    slide = prs.slides[0]
    for shape in slide.shapes:
        print(shape.has_table)

    with open(csv_path,'r',encoding='utf-8_sig') as f:
        reader = csv.reader(f)
        sample_data = [row for row in reader]

    # 評価星リストの作成
    num = re.findall('[0-9]+',sample_data[3][1])
    star_list = list(map(float, num)) # 文字列配列を数値配列に
    star_list = star_list[1::2] #奇数要素を取得
    for i in range(len(star_list)):
        star_list[i]=star_list[i]/100

    # グラフの書き換え
    chart = prs.slides[0].shapes[2].chart  
    chart_data = CategoryChartData()
    chart_data.categories = '星5つ','星4つ','星3つ','星2つ','星1つ'
    chart_data.add_series('系列 1', star_list)
    chart.replace_data(chart_data)

    # csvファイルからの書き換え
    font_size = Pt(18)
    title_font_size = Pt(20)
    font_color = RGBColor(0xFF, 0xFF, 0xFF)
    line_spacing = 1.4
    with open(csv_path,'r',encoding='utf-8_sig') as f:
        reader = csv.reader(f)
        sample_data = [row for row in reader]
    for i in range(len(sample_data)):
        if i < 4 :
            pass
        elif 4 <= i <= 11:
            tbl = prs.slides[0].shapes[math.floor(i/2)+1].table
            # 偶数行目に白文字
            for j in range(len(sample_data[i])):
                tbl.cell(i%2, j).text =str(sample_data[i][j])
                paras = tbl.cell(i%2, j).text_frame.paragraphs
                for para in paras:
                    para.font.size = font_size
                    para.font.color.rgb = font_color
                    if j % 2 == 1:
                        para.font.size = title_font_size
                        para.font.bold = True
                    if i % 2 == 1:
                        para.font.color.rgb = font_color
                        para.line_spacing = line_spacing
    dst_path = f'img/{date_now}.png' 
    width = prs.slide_width
    height = prs.slide_height
    pic = prs.slides[0].shapes.add_picture(dst_path, 0, 0, width = width, height = None)
    if height - pic.height < 0 :
        del_pic = prs.slides[0].shapes[7] 
        del_pic = del_pic._element
        del_pic.getparent().remove(del_pic)
        pic = prs.slides[0].shapes.add_picture(dst_path, 0, 0, width = None, height = height)
    # 最下部中心に移動
    pic.left = int( ( width  - pic.width  ) / 2 )
    pic.top  = int( height - pic.height )
    # 重なりの順序を変更
    picture = prs.slides[0].shapes[0]._element 
    picture.addprevious(prs.slides[0].shapes[7]._element) # shape[2]をshape[0]の前つまり最背面に移動
    prs.save(output_path)

def do():
    country = 'JAPAN'
    #商品ページ
    url = 'https://www.amazon.co.jp/GP-AOTGP-CTNS6042-カタン-スタンダード版/dp/B017SB7QLO/?_encoding=UTF8&pd_rd_w=o0xEW&pf_rd_p=d753a891-bba8-4352-b10c-3c8d00c6b548&pf_rd_r=V6EC7AT60W8PF3ZZ2WX2&pd_rd_r=af2ce670-afe9-4fbf-be88-b47718ed2f87&pd_rd_wg=FbtZK&ref_=pd_gw_ci_mcx_mr_hp_atf_m&th=1'
    url= url.encode('ascii', 'ignore').decode('unicode_escape')
    review_url = url.replace('dp', 'product-reviews')
    date_now = write_csv(url,review_url)
    write_pptx(country,date_now)
    create_short(date_now)

if __name__ == '__main__':
    do()