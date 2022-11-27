from selenium import webdriver
import chromedriver_binary
import time
from time import sleep
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.common.by import By
from bs4 import BeautifulSoup
import pandas as pd
import requests
import json
import textwrap
import csv
import datetime
import pptx
import urllib.error
import urllib.request
from pptx.util import Pt
import shutil
import math
import sys

URL = 'https://www.googleapis.com/youtube/v3/'
API_KEY = 'AIzaSyBXlTVtjMZI3eeJZaByEHPlMmWJvBCbvXU'



def get_page_info(url):
    text = ""                              
    options = Options()                    
    options.add_argument('--incognito')    
    #options.add_argument('--headless')      
    driver = webdriver.Chrome(options=options)
    driver.get(url)                         
    driver.implicitly_wait(10)         
    text = driver.page_source
    # 最下部までスクロールして更新されたらtextを変更、変化なしの場合スクロール終了break
    while 1:
        driver.find_element(By.TAG_NAME,'body').send_keys(Keys.END)
        sleep(3)
        text2=driver.page_source
        if text!=text2:
            text=text2
        else:
            break
    driver.quit()                           
    return text                             

def get_status(url):
    id_list = []
    text = get_page_info(url)    
    bs = BeautifulSoup(text, features='lxml')  
    title =  bs.select('a#video-title')
    for i in range(len(title)):
        id_elem = []
        video_id = bs.select('a#video-title')[i].get('href')
        title =  bs.select('a#video-title')[i].text
        if video_id :
            video_id = video_id[-11:]
            video_url = f'https://youtu.be/{video_id}'
            id_elem.append(video_id)
            id_elem.append(video_url)
            id_elem.append(title)
            id_list.append(id_elem)
    return id_list

# video_id_listのcsv出力
def write_csv(url,csv_path):
    id_list = get_status(url)
    df = pd.DataFrame(id_list,columns=['video_id','url','title'])
    df.to_csv(csv_path, index=False)


def print_video_comment(no, video_id, next_page_token, text_data):
    params = {
      'key': API_KEY,
      'part': 'snippet',
      'videoId': video_id,
      # order=timeにすると上限なしでコメントを取得可能らしい
      # order=relevanceで上限2000
      'order': 'time',
      'textFormat': 'plaintext',
      'maxResults': 100,
    }
    if next_page_token is not None:
        params['pageToken'] = next_page_token
    response = requests.get(URL + 'commentThreads', params=params)
    resource = response.json()

    for comment_info in resource['items']:
        # コメント
        text = comment_info['snippet']['topLevelComment']['snippet']['textDisplay']
        # グッド数
        like_cnt = comment_info['snippet']['topLevelComment']['snippet']['likeCount']
        # 返信数
        reply_cnt = comment_info['snippet']['totalReplyCount']
        # ユーザー名
        user_name = comment_info['snippet']['topLevelComment']['snippet']['authorDisplayName']
        # Id
        parentId = comment_info['snippet']['topLevelComment']['id']
        # ユーザープロフィールimage
        author_channel = comment_info["snippet"]['topLevelComment']['snippet']['authorProfileImageUrl']
        # 最終更新日時
        published_date = comment_info["snippet"]['topLevelComment']['snippet']['publishedAt']
        # リストに取得した情報を追加
        text_data.append([parentId, 'parent', text, like_cnt, reply_cnt, user_name, author_channel, published_date])
        # 処理確認用
        if len(text_data) % 100 == 0:
            print(len(text_data))
#     print('{:0=4}\t{}\t{}\t{}\t{}'.format(no, text.replace('\n', ' '), like_cnt, user_name, reply_cnt))
        no = no + 1
    if 'nextPageToken' in resource:
        print_video_comment(no, video_id, resource["nextPageToken"], text_data)
        

# download img from web
def download_file(url, dst_path):
    try:
        with urllib.request.urlopen(url) as web_file:
            data = web_file.read()
            with open(dst_path, mode='wb') as local_file:
                local_file.write(data)
    except urllib.error.URLError as e:
        print(e)


# shorts動画作成
def create_short(video_id):
    csv_path=f'csv/{video_id}.csv'
    output_path=f'pptx/short_{video_id}.pptx'
    src = 'pptx/short.pptx'
    copy = output_path
    shutil.copyfile(src,copy)

    prs = pptx.Presentation(output_path)

    # スライド内のオブジェクトの読み取り
    # slide = prs.slides[0]
    # for shape in slide.shapes:
    #     print(shape.has_table) #  False True true true
    
    df = pd.read_csv(csv_path, sep=',', engine='python')
    top_num = 3
    top_elem = ['like_cnt','user_name','comment_data','published_date'] 
    top_data = df[top_elem].head(top_num)
    font_size = Pt(20)
    for i in range(len(top_data)):
        tbl = prs.slides[0].shapes[i+1].table
        tbl.cell(0, 0).text = '👍🏿 '+str(top_data.at[i,'like_cnt'])
        paras = tbl.cell(0, 0).text_frame.paragraphs
        for para in paras:
            para.font.size = font_size
            para.font.bold = True
        date = datetime.datetime.fromisoformat(top_data.at[i,'published_date'].replace('Z', ''))
        pass_date = calc_date(date)
        tbl.cell(0, 1).text = top_data.at[i,'user_name']+'      '+pass_date
        tbl.cell(1, 0).text = top_data.at[i,'comment_data']
        paras = tbl.cell(1, 0).text_frame.paragraphs
        for para in paras:
            para.font.size = font_size
            para.font.bold = True
    prs.save(output_path)

    # csvファイルからの書き換え
    # font_size = Pt(18)
    # title_font_size = Pt(20)
    # font_color = RGBColor(0xFF, 0xFF, 0xFF)
    # line_spacing = 1.4

    # para.font.bold = True
    # para.line_spacing = line_spacing

    dst_path = f'img/{video_id}.png' 
    width = prs.slide_width
    height = prs.slide_height
    pic = prs.slides[0].shapes.add_picture(dst_path, 0, 0, width = width, height = None)
    # 最上部中心に移動
    pic.left = int( 0 )
    pic.top  = int( 0 )
    # 重なりの順序を変更
    picture = prs.slides[0].shapes[0]._element 
    picture.addprevious(prs.slides[0].shapes[5]._element) 
    prs.save(output_path)


# youtube_shorts動画作成
def create_youtube_short(video_id):
    csv_path=f'csv/{video_id}.csv'
    output_path=f'pptx/youtube_short_{video_id}.pptx'
    src = 'pptx/youtube_short.pptx'
    copy = output_path
    shutil.copyfile(src,copy)

    prs = pptx.Presentation(output_path)

    # スライド内のオブジェクトの読み取り
    # slide = prs.slides[0]
    # for shape in slide.shapes:
    #     print(shape.has_table) #  False True true true
    
    df = pd.read_csv(csv_path, sep=',',engine='python')
    top_num = 3
    top_elem = ['like_cnt','user_name','comment_data','published_date'] 
    top_data = df[top_elem].head(top_num)
    font_size = Pt(20)
    for i in range(len(top_data)):
        tbl = prs.slides[0].shapes[i+2].table
        tbl.cell(0, 0).text = '👍🏿 '+str(top_data.at[i,'like_cnt'])
        paras = tbl.cell(0, 0).text_frame.paragraphs
        for para in paras:
            para.font.size = font_size
            para.font.bold = True
        date = datetime.datetime.fromisoformat(top_data.at[i,'published_date'].replace('Z', ''))
        pass_date = calc_date(date)
        tbl.cell(0, 1).text = top_data.at[i,'user_name']+'      '+pass_date
        tbl.cell(1, 0).text = top_data.at[i,'comment_data']
        paras = tbl.cell(1, 0).text_frame.paragraphs
        for para in paras:
            para.font.size = font_size
            para.font.bold = True
    prs.save(output_path)

    # csvファイルからの書き換え
    # font_size = Pt(18)
    # title_font_size = Pt(20)
    # font_color = RGBColor(0xFF, 0xFF, 0xFF)
    # line_spacing = 1.4

    # para.font.bold = True
    # para.line_spacing = line_spacing

    dst_path = f'img/{video_id}.png' 
    width = prs.slide_width
    height = prs.slide_height
    pic = prs.slides[0].shapes.add_picture(dst_path, 0, 0, width = width, height = None)
    # 最上部中心に移動
    pic.left = int( 0 )
    pic.top  = int( 0 )
    # 重なりの順序を変更
    picture = prs.slides[0].shapes[0]._element 
    picture.addprevious(prs.slides[0].shapes[5]._element) 
    prs.save(output_path)


# 日付変換メソッド
def calc_date(date):
    today = datetime.datetime.today()
    pass_date = today - date
    pass_date = pass_date.days
    year = pass_date // 365
    month = pass_date // 30
    week = pass_date // 7
    if year > 0 :
        return f'{year}年前'
    elif month > 0 :
        return f'{month}ヶ月前'
    elif week > 0 :
        return f'{week}週間前'
    else :
        return f'{pass_date}日前'


# 1minute動画作成
def create_1minute(video_id):
    csv_path=f'csv/{video_id}.csv'
    output_path=f'pptx/{video_id}.pptx'
    src = 'pptx/1minute.pptx'
    copy = output_path
    shutil.copyfile(src,copy)

    prs = pptx.Presentation(output_path)

    # スライド内のオブジェクトの読み取り
    # slide = prs.slides[0]
    # for shape in slide.shapes:
    #     print(shape.has_table) #  False True true true
    
    df = pd.read_csv(csv_path, sep=',',engine='python')
    top_num = 10
    top_elem = ['like_cnt','user_name','comment_data','published_date'] 
    top_data = df[top_elem].head(top_num)

    font_size = Pt(28)
    for i in range(len(top_data)):
        tbl = prs.slides[0].shapes[11-i].table
        tbl.cell(0, 0).text = '👍🏿 '+str(top_data.at[i,'like_cnt'])
        paras = tbl.cell(0, 0).text_frame.paragraphs
        for para in paras:
            para.font.size = font_size
            para.font.bold = True
        date = datetime.datetime.fromisoformat(top_data.at[i,'published_date'].replace('Z', ''))
        pass_date = calc_date(date)
        tbl.cell(0, 1).text = top_data.at[i,'user_name']+'      '+pass_date
        tbl.cell(1, 0).text = top_data.at[i,'comment_data']
        paras = tbl.cell(1, 0).text_frame.paragraphs
        for para in paras:
            para.font.size = font_size
            para.font.bold = True

    # csvファイルからの書き換え
    # font_size = Pt(18)
    # title_font_size = Pt(20)
    # font_color = RGBColor(0xFF, 0xFF, 0xFF)
    # line_spacing = 1.4

    # para.font.bold = True
    # para.line_spacing = line_spacing

    dst_path = f'img/{video_id}.png' 
    width = prs.slide_width
    height = prs.slide_height
    pic = prs.slides[0].shapes.add_picture(dst_path, 0, 0, width = None, height = height)
    # 最上部中心に移動
    pic.left = int( ( width  - pic.width  ) / 2 )
    pic.top  = int( 0 )
    # 重なりの順序を変更
    picture = prs.slides[0].shapes[0]._element 
    picture.addprevious(prs.slides[0].shapes[12]._element) 
    prs.save(output_path)


def do():

    # 動画の一覧ページ
    url = 'https://www.youtube.com/user/TokaiOnAir/videos?view=0&sort=p&shelf_id=0'
    # videoidのリスト
    video_list_name = '東海オンエア　人気'


    csv_path = 'csv/'+str(video_list_name)+'.csv'
    #write_csv(url,csv_path)
    df = pd.read_csv(csv_path, sep=',',engine='python')
    # videoidをid_listに格納
    id_list = df[['video_id']]
    
    for id in range(171,202):
        video_id = id_list.at[id,'video_id']
        text_data=[]
        # コメントを全取得する

        no = 1
        # 取得処理を実行
        print_video_comment(no, video_id, None, text_data)
        # データフレーム作成(高評価順にソート)
        df = pd.DataFrame(text_data, columns=['comment_id', 'type', 'comment_data', 'like_cnt', 'reply_cnt', 'user_name', 'profile_page', 'published_date']).sort_values('like_cnt', ascending=False)
        # csv出力
        df.to_csv('csv/'+str(video_id)+'.csv', index=False)
        # データフレームを一部出力して確認する
        df.head()
        print(f'{id} comment csv finished!')    

        url = f'http://img.youtube.com/vi/{video_id}/hqdefault.jpg'
        dst_path = f'img/{video_id}.png'
        download_file(url, dst_path)
        #create_short(video_id)
        #print(f'{id} short-pptx finished!') 
        
        create_youtube_short(video_id)
        print(f'{id} youtube_short-pptx finished!') 
        
        create_1minute(video_id)
        print(f'{id} pptx-1minute finished!') 


if __name__ == '__main__':
    do()